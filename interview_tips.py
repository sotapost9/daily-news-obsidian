import sys, os
_SKILL = os.path.join(os.path.expanduser("~"), ".claude", "skills", "acnpptx", "scripts")
sys.path.insert(0, _SKILL)

import helpers as _h
_h.set_lang("ja")
_h.load_theme("acn-slide-master")
from helpers import *

from native_shapes import *
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn as _qn

prs = Presentation(TEMPLATE_PATH)
while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    prs.part.drop_rel(sldId.get(_qn("r:id")))
    del prs.slides._sldIdLst[0]

LAYOUT_COVER = 0
LAYOUT_CONTENT = 57


def add_slide(layout_idx=LAYOUT_CONTENT):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    clear_placeholders(slide)
    return slide


def add_cover_slide():
    return prs.slides.add_slide(prs.slide_layouts[LAYOUT_COVER])


# ─── Slide 1: Cover ──────────────────────────────────────────────────────────
slide = add_cover_slide()
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    ph.text_frame.clear()
    p = ph.text_frame.paragraphs[0]
    if idx == 0:
        p.text = "インタビューのコツ"
        p.font.size = Pt(40); p.font.bold = True
        p.font.color.rgb = WHITE; p.font.name = FONT
    elif idx == 1:
        p.text = "Sotaの個人Vaultから抽出した実践ナレッジ"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE; p.font.name = FONT
    elif idx == 12:
        p.text = "Sota Yamanaka"
        p.font.size = Pt(14); p.font.color.rgb = LIGHT_PURPLE; p.font.name = FONT
    elif idx == 15:
        p.text = "2026年6月"
        p.font.size = Pt(8); p.font.color.rgb = LIGHT_PURPLE; p.font.name = FONT
    else:
        p.text = " "; p.font.size = Pt(8); p.font.name = FONT


# ─── Slide 2: Agenda (Pattern I) ─────────────────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "アジェンダ")
add_title(slide, "アジェンダ")

agenda_items = [
    "マインドセット",
    "事前準備",
    "場の作り方",
    "質問の技術",
    "聞き方の姿勢",
    "深掘りの型（ジョブ理論）",
    "インタビュー後の整理",
]
ITEM_Y = CY + 0.10
gap = min(0.78, (BY - ITEM_Y - 0.65) / max(len(agenda_items) - 1, 1))

for i, item in enumerate(agenda_items):
    iy = ITEM_Y + i * gap
    num_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(ML), Inches(iy), Inches(0.55), Inches(0.62))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = DARKEST_PURPLE
    num_box.line.fill.background()
    tf = num_box.text_frame
    tf.word_wrap = False
    _bp = tf._txBody.find(_qn("a:bodyPr"))
    for attr in ("lIns", "rIns", "tIns", "bIns"):
        _bp.set(attr, "0")
    _bp.set("anchor", "ctr")
    tf.paragraphs[0].text = f"{i+1:02d}"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = FONT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(
        Inches(ML + 0.70), Inches(iy + 0.08), Inches(11.30), Inches(0.50))
    p = tb.text_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(18); p.font.color.rgb = TEXT_BODY; p.font.name = FONT

set_footer(slide)


# ─── Slide 3: マインドセット (Pattern D) ─────────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "マインドセット")
add_title(slide, "マインドセット：全ての土台")
add_message_line(slide, "テクニックより「相手への純粋な関心」がインタビューの成否を決める")

# Accent bar
bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    Inches(ML), Inches(2.30), Inches(1.20), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = CORE_PURPLE; bar.line.fill.background()

# Key message
tb_key = slide.shapes.add_textbox(Inches(ML), Inches(2.45), Inches(CW), Inches(1.60))
tf = tb_key.text_frame; tf.word_wrap = True
tf._txBody.find(_qn("a:bodyPr")).set("anchor", "t")
p = tf.paragraphs[0]
p.text = "どんなに巧みな質問技術より、「相手を知りたい」という純粋な関心が重要。\n関心があれば、自然と「その前後で何があったか」「その時どう感じたか」という問いが生まれる。"
p.font.size = Pt(28); p.font.bold = True
p.font.color.rgb = BLACK; p.font.name = FONT
p.space_after = Pt(10)

# Supporting bullets
support_items = [
    "「映画監督」になる：抽象的な意見ではなく、具体的な過去の行動と状況を映像として捉える",
    "目的から逆算する：「何を特定・検証するために、何を理解するか」を事前に言語化する",
    "2つの視点を持つ：没入する自分と、「聞くべきことを聞けているか」チェックするもう一人の自分",
    "打算的な下心は見抜かれる：「聞き出してやろう」という態度は相手に伝わり、場が閉じる",
]

tb_sup = slide.shapes.add_textbox(Inches(ML), Inches(4.20), Inches(CW), Inches(BY - 4.20 - 0.10))
tf = tb_sup.text_frame; tf.word_wrap = True
tf._txBody.find(_qn("a:bodyPr")).set("anchor", "t")
for j, line in enumerate(support_items):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.text = f"• {line}"; p.font.size = Pt(14)
    p.font.color.rgb = TEXT_BODY; p.font.name = FONT
    p.space_after = Pt(8)

set_footer(slide)


# ─── Slide 4: 事前準備 (Pattern P) ───────────────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "事前準備")
add_title(slide, "事前準備の流れ")
add_message_line(slide, "目的から逆算した設計と事前調査が、場での余裕と深掘りを生む")

steps = ["①目的の定義", "②対象者の調査", "③フロー設計", "④練習"]
details = [
    "• 「何を特定・検証するためにYYを理解する」と言語化\n• NG: 「ニーズを知る」\n• OK: 「30代男性が休日の朝にコーヒーを飲む際の情緒的価値を特定する」",
    "• 相手のSNS・記事・プロフィールを徹底調査\n• 事前調査が場での余裕を生み、深掘りにつながる\n• N1（実在する一人）を深掘りする意識を持つ",
    "• 一言一句の台本ではなくチェックリスト形式\n• 導入→本編（広い質問→具体）→結びの3段構成\n• 目的・仮説のないインタビューはチグハグになる",
    "• 実際に質問を声に出して練習する\n• 鏡の前や信頼できる人との模擬インタビュー\n• 質問の流れとタイミングを体に染み込ませる",
]

FLOW_H = 0.75
GAP = 0.20
DETAIL_H = 3.00
FLOW_Y = CY + (AH - FLOW_H - GAP - DETAIL_H) / 2
DETAIL_Y = FLOW_Y + FLOW_H + GAP

add_chevron_flow(
    slide, steps,
    x=ML, y=FLOW_Y, total_w=CW, h=FLOW_H,
    fill_color=DARKEST_PURPLE,
    text_color=WHITE,
    shape_style="chevron",
    use_pentagon_first=True,
    font_size_pt=14,
    font_name=FONT,
)

n = len(steps)
col_w = CW / n
for i, detail in enumerate(details):
    dx = ML + i * col_w
    tb = slide.shapes.add_textbox(Inches(dx + 0.08), Inches(DETAIL_Y),
        Inches(col_w - 0.16), Inches(DETAIL_H))
    tf = tb.text_frame; tf.word_wrap = True
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "t")
    for j, line in enumerate(detail.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(12)
        p.font.color.rgb = TEXT_BODY; p.font.name = FONT
        p.space_after = Pt(7)

set_footer(slide)


# ─── Slide 5: 場の作り方 (Pattern K) ─────────────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "場の作り方")
add_title(slide, "場の作り方")
add_message_line(slide, "インタビューの質は、最初の数分間でほぼ決まる")

N_COLS = 3
PANEL_W = (CW - 0.20 * (N_COLS - 1)) / N_COLS
PANEL_H = BY - CY - 0.60
BADGE_D = 0.48

cols = [
    {
        "num": 1, "title": "ラポール形成",
        "body": "• 相手へのリスペクトを言葉と態度で伝える\n• 相手のSNSや実績に「本物の関心」を示す\n• 自分のことも少し話してフラットな関係を作る\n• 「先にジャブを打つ」：共通点や話しやすい話題を先に出す\n• ラポールなき深掘りは、尋問になる"
    },
    {
        "num": 2, "title": "流れを最初に説明する",
        "body": "• 「今日は〇〇について約60分お聞きします」と全体像を伝える\n• ゴールと進め方を共有すると相手が安心して話せる\n• 録音・録画の可否を最初に確認する\n• 守秘義務・使用目的を明示して信頼を担保する\n• 「最後に何でも聞いてください」と伝えて双方向感を出す"
    },
    {
        "num": 3, "title": "空気を作る",
        "body": "• インタビュアーが場の空気をコントロールする意識を持つ\n• 相手がどんな状態で入ってくるかは最初で決まる\n• 自分が落ち着いていれば、相手も落ち着く\n• 緊張している時こそ、ゆっくり話す・間を怖がらない\n• 「正解はない」「思ったことを率直に」と念押しする"
    },
]

for i, col in enumerate(cols):
    px = ML + i * (PANEL_W + 0.20)
    py = CY + 0.55

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(px), Inches(py), Inches(PANEL_W), Inches(PANEL_H))
    bg.fill.solid(); bg.fill.fore_color.rgb = OFF_WHITE; bg.line.fill.background()

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(px), Inches(py), Inches(PANEL_W), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = CORE_PURPLE; bar.line.fill.background()

    bx = px + (PANEL_W - BADGE_D) / 2
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(bx), Inches(py - BADGE_D / 2), Inches(BADGE_D), Inches(BADGE_D))
    badge.fill.solid(); badge.fill.fore_color.rgb = DARKEST_PURPLE; badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = str(col["num"])
    p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "ctr")

    tb_t = slide.shapes.add_textbox(Inches(px + 0.12), Inches(py + 0.16),
        Inches(PANEL_W - 0.24), Inches(0.50))
    p = tb_t.text_frame.paragraphs[0]
    p.text = col["title"]
    p.font.size = Pt(14); p.font.bold = True
    p.font.color.rgb = BLACK; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER

    tb_b = slide.shapes.add_textbox(Inches(px + 0.12), Inches(py + 0.75),
        Inches(PANEL_W - 0.24), Inches(PANEL_H - 0.90))
    tf = tb_b.text_frame; tf.word_wrap = True
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "t")
    for j, line in enumerate(col["body"].split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(12)
        p.font.color.rgb = TEXT_BODY; p.font.name = FONT
        p.space_after = Pt(5)

set_footer(slide)


# ─── Slide 6: 質問の技術 (Pattern F 2×2) ─────────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "質問の技術")
add_title(slide, "質問の技術")
add_message_line(slide, "過去の具体的な行動から聞き、仮説を持って広い質問→具体へ深める")

CARD_W = (CW - 0.30) / 2
CARD_H = (BY - CY - 0.35 - 0.20) / 2
xs = [ML, ML + CARD_W + 0.30]
ys_card = [CY + 0.20, CY + 0.20 + CARD_H + 0.20]

cards = [
    {
        "title": "過去の行動を聞く",
        "body": "• 「最後にXXしたのはいつですか？」から開始\n• 人間は未来を予測できないが、過去の行動は事実\n• 「あったら使いますか？」は聞かない\n• 日時を特定すると状況・感情・環境が芋づる式に出てくる",
    },
    {
        "title": "広い質問→具体へ",
        "body": "• まずエピソード記憶を呼び起こす広い質問から\n• その後「いつ・どこで・誰と・何を」で状況を映画のシーンのように描写させる\n• 抽象→具体のリズムを意識してセッションを進める",
    },
    {
        "title": "仮説を当てにいく",
        "body": "• 相手がうまく言語化できない時に「こういうことないですか？」と仮説を提示\n• 事前に仮説を持っておくことが前提\n• 押しつけではなく「確認」として出すのがコツ\n• 仮説が外れた時こそ深いインサイトが生まれる",
    },
    {
        "title": "クローズド質問に注意",
        "body": "• 相手の語りたがり度が低いとYES/NO回答になりやすい\n• 「〜ですよね？」「〜でしたか？」はクローズド\n• 「〜について教えてもらえますか？」に変える習慣\n• 誘導尋問は分析に使えないデータを生む",
    },
]

for idx_c, card in enumerate(cards):
    px = xs[idx_c % 2]
    py = ys_card[idx_c // 2]

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(px), Inches(py), Inches(CARD_W), Inches(CARD_H))
    bg.fill.solid(); bg.fill.fore_color.rgb = OFF_WHITE; bg.line.fill.background()

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(px), Inches(py), Inches(CARD_W), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = CORE_PURPLE; bar.line.fill.background()

    tb_h = slide.shapes.add_textbox(Inches(px + 0.15), Inches(py + 0.12),
        Inches(CARD_W - 0.30), Inches(0.40))
    p = tb_h.text_frame.paragraphs[0]
    p.text = card["title"]
    p.font.size = Pt(14); p.font.bold = True
    p.font.color.rgb = BLACK; p.font.name = FONT

    tb_b = slide.shapes.add_textbox(Inches(px + 0.15), Inches(py + 0.58),
        Inches(CARD_W - 0.30), Inches(CARD_H - 0.68))
    tf = tb_b.text_frame; tf.word_wrap = True
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "t")
    for j, line in enumerate(card["body"].split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(12)
        p.font.color.rgb = TEXT_BODY; p.font.name = FONT
        p.space_after = Pt(6)

set_footer(slide)


# ─── Slide 7: 聞き方の姿勢 (Pattern E) ───────────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "聞き方の姿勢")
add_title(slide, "聞き方の姿勢")
add_message_line(slide, "発言を素直に受け取り、インタビュー中に結論を出そうとしない")

e_items = [
    (
        "誘導しない",
        "分析に使えるのはユーザーの発言だけ。自分の仮説を押しつけない。「XXだから便利ですよね？」はNG。報奨金のような前提を置いた質問も、相手が気を遣っているだけになる。",
    ),
    (
        "言葉をそのまま受け止める",
        "自分の解釈を入れない。発言は事実として、私見と切り分ける。「つまりこういうことですか？」と確認はしてよいが、自分の主観的解釈を前提にしてはいけない。その場で生まれる言葉を大切にし、「このセリフを言って欲しい」という下心を出さない。",
    ),
    (
        "結論を急がない",
        "インタビュー中に結論を出そうとしない。目の前のN1の事実は事実として受け止め、分析は後で行う。確からしい示唆が出ても、他のインタビュイーは全く異なる行動をしている可能性がある。全体の矛盾を見つける意識を持つ。",
    ),
]

row_h = (BY - CY - 0.30) / len(e_items)
for i, (headline, detail) in enumerate(e_items):
    row_y = CY + 0.15 + i * row_h

    acc = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(ML), Inches(row_y + 0.06), Inches(0.06), Inches(0.30))
    acc.fill.solid(); acc.fill.fore_color.rgb = CORE_PURPLE; acc.line.fill.background()

    tb_h = slide.shapes.add_textbox(Inches(ML + 0.20), Inches(row_y),
        Inches(CW - 0.20), Inches(0.42))
    p = tb_h.text_frame.paragraphs[0]
    p.text = headline; p.font.size = Pt(18)
    p.font.bold = True; p.font.color.rgb = TEXT_BODY; p.font.name = FONT

    tb_d = slide.shapes.add_textbox(Inches(ML + 0.20), Inches(row_y + 0.45),
        Inches(CW - 0.20), Inches(row_h - 0.55))
    tf = tb_d.text_frame; tf.word_wrap = True
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "t")
    p = tf.paragraphs[0]
    p.text = detail; p.font.size = Pt(14)
    p.font.color.rgb = TEXT_BODY; p.font.name = FONT
    p.space_after = Pt(8)

set_footer(slide)


# ─── Slide 8: 深掘りの型 (Pattern G) ─────────────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "深掘りの型")
add_title(slide, "深掘りの型（ジョブ理論ベース）")
add_message_line(slide, "5つの問いで「本質的なジョブ」を表面ニーズの奥から引き出す")

headers = ["ステップ", "問いの例", "目的", "ポイント"]
rows_data = [
    ["① Job（進歩）",    "「その時、何を達成しようとしていましたか？」",      "本質的な目的の特定",   "「〜したい」という動詞で捉える"],
    ["② Context（状況）","「何時頃でしたか？誰といましたか？どこで？」",       "行動の文脈・環境の把握", "時間・場所・同伴者で状況を具体化"],
    ["③ Barrier（障害）","「何が邪魔をしていましたか？うまくいかない原因は？」","阻害要因の発見",       "「不満」ではなく「何が邪魔か」で聞く"],
    ["④ Solution（対処）","「今はどう対処していますか？何が不満ですか？」",    "現在の代替行動と不満点","不完全な解決策の中にニーズが潜む"],
    ["⑤ Criteria（条件）","「理想の解決策は？譲れない条件は何ですか？」",      "意思決定の基準の把握",  "「あれば嬉しい」と「必須」を分ける"],
]

tbl_rows = len(rows_data) + 1
tbl_cols = len(headers)
table = slide.shapes.add_table(
    tbl_rows, tbl_cols,
    Inches(ML), Inches(2.05), Inches(CW), Inches(BY - 2.15)
).table

col_widths = [1.80, 3.80, 2.80, 4.10]
for c, w in enumerate(col_widths):
    table.columns[c].width = Inches(w)

for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = DARKEST_PURPLE
    for para in cell.text_frame.paragraphs:
        para.font.bold = True; para.font.size = Pt(12)
        para.font.color.rgb = WHITE; para.font.name = FONT
        para.alignment = PP_ALIGN.CENTER

for r, row in enumerate(rows_data, 1):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(12)
            if c == 0:
                para.font.bold = True; para.font.color.rgb = DARKEST_PURPLE
            else:
                para.font.color.rgb = TEXT_BODY
            para.font.name = FONT

set_footer(slide)


# ─── Slide 9: インタビュー後の整理 (Pattern T) ───────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "インタビュー後")
add_title(slide, "インタビュー後の整理")
add_message_line(slide, "ラップアップとJOBS分解で、生の発言をインサイトへ昇華させる")

LABEL_W = 2.00
ARROW_H = 0.45
GAP = 0.15
section_h = (BY - CY - 0.10 - ARROW_H - GAP * 3) / 2

sections_data = [
    {
        "label": "その場で\nやること",
        "items": [
            "• ラップアップ（認識合わせ）：「今日お話いただいた中で、XXXとおっしゃっていましたが、つまりYYYということですか？」と確認",
            "• 言い残したことの確認：「他に何かあれば教えてください」で追加情報を引き出す",
            "• 矛盾・引っかかりのメモ：インタビュー中に気になった点を手元にメモしておく",
        ]
    },
    {
        "label": "後処理で\nやること",
        "items": [
            "• JOBS分解：発言を付箋に書き出し、Job / Objective / Barrier / Solution の4要素に分類する",
            "• 3つのジョブで分析：機能的ジョブ（タスク）・感情的ジョブ（どう感じたいか）・社会的ジョブ（どう見られたいか）",
            "• ジョブスペックの定義：製品仕様ではなく「体験の性能条件」として整理する（例：10分で終わる・予約不要・清潔感がある）",
        ]
    },
]

for i, sec in enumerate(sections_data):
    sy = CY + 0.05 + i * (section_h + ARROW_H + GAP)

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(ML), Inches(sy), Inches(CW), Inches(section_h))
    bg.fill.solid(); bg.fill.fore_color.rgb = OFF_WHITE
    bg.line.color.rgb = LIGHT_GRAY; bg.line.width = Pt(0.75)

    lbl = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(ML), Inches(sy), Inches(LABEL_W), Inches(section_h))
    lbl.fill.solid(); lbl.fill.fore_color.rgb = DARKEST_PURPLE; lbl.line.fill.background()
    tf = lbl.text_frame; tf.word_wrap = True
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "ctr")
    p = tf.paragraphs[0]; p.text = sec["label"]
    p.font.size = Pt(14); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER

    content_x = ML + LABEL_W + 0.15
    content_w = CW - LABEL_W - 0.20
    tb = slide.shapes.add_textbox(Inches(content_x), Inches(sy + 0.12),
        Inches(content_w), Inches(section_h - 0.25))
    tf = tb.text_frame; tf.word_wrap = True
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "ctr")
    for j, line in enumerate(sec["items"]):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(13)
        p.font.color.rgb = TEXT_BODY; p.font.name = FONT
        p.space_after = Pt(10)

    if i == 0:
        arrow_y = sy + section_h + GAP
        arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
            Inches(ML + CW / 2 - 0.50), Inches(arrow_y), Inches(1.0), Inches(ARROW_H))
        arr.fill.solid(); arr.fill.fore_color.rgb = DARKEST_PURPLE; arr.line.fill.background()

set_footer(slide)


# ─── Closing ──────────────────────────────────────────────────────────────────
make_closing_slide(prs)
strip_sections(prs)

output_path = os.path.join(os.getcwd(), "インタビューのコツ.pptx")
prs.save(output_path)
print(f"✅ 保存完了: {output_path}")
