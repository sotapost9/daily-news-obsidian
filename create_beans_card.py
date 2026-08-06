#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys, os
_SKILL = os.path.join(os.path.expanduser("~"), ".claude", "skills", "acnpptx", "scripts")
sys.path.insert(0, _SKILL)

import helpers as _h
_h.set_lang("ja")
_h.load_theme("accenture")
from helpers import *

from native_shapes import *
from charts import add_bar_chart
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn as _qn

prs = Presentation(TEMPLATE_PATH)
while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    prs.part.drop_rel(sldId.get(_qn("r:id")))
    del prs.slides._sldIdLst[0]

# ===== COVER SLIDE =====
slide = prs.slides.add_slide(prs.slide_layouts[0])
clear_placeholders(slide)

# プレースホルダー埋め込み
ph_map = {}
for idx, placeholder in enumerate(slide.placeholders):
    ph_map[idx] = placeholder

# タイトル
if 0 in ph_map:
    ph_map[0].text = "Accenture\nCoffee Beans\nTasting Card"
    p = ph_map[0].text_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(88, 24, 112)  # Accenture Purple

if 1 in ph_map:
    ph_map[1].text = "社内 Fes"
    p = ph_map[1].text_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(88, 24, 112)

if 2 in ph_map:
    ph_map[2].text = "2026年8月"
    p = ph_map[2].text_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(100, 100, 100)

set_footer(slide)

# ===== SLIDE 2: LIGHT ROAST (レーダーチャート) =====
slide2 = prs.slides.add_slide(prs.slide_layouts[2])
clear_placeholders(slide2)

add_breadcrumb(slide2, "Coffee Beans > Light Roast")
add_title(slide2, "Light Roast（浅煎り）")
add_message_line(slide2, "爽やかな酸味と明るい香りが特徴。フルーティーな風味が引き立つ")

# レーダーチャート用データ
categories = ["酸味", "コク", "甘さ", "苦味", "香り"]
light_roast_values = [4, 2, 3, 1, 4]

# 横棒グラフで表現（レーダーは複雑なので）
add_bar_chart(slide2, title="テイスティングプロファイル",
              categories=categories,
              series_data=[{"name": "評価", "values": light_roast_values}],
              x=ML, y=Inches(1.80), w=CW, h=Inches(3.80),
              show_data_labels=True, font_name=FONT)

set_footer(slide2)

# ===== SLIDE 3: MEDIUM ROAST (横棒グラフ) =====
slide3 = prs.slides.add_slide(prs.slide_layouts[2])
clear_placeholders(slide3)

add_breadcrumb(slide3, "Coffee Beans > Medium Roast")
add_title(slide3, "Medium Roast（中煎り）")
add_message_line(slide3, "バランスの取れた熱い味わい。豊かなボディとクリーンな後味が特徴")

categories = ["酸味", "コク", "甘さ", "苦味", "香り"]
medium_roast_values = [3, 4, 4, 2, 4]

add_bar_chart(slide3, title="テイスティングプロファイル",
              categories=categories,
              series_data=[{"name": "評価", "values": medium_roast_values}],
              x=ML, y=Inches(1.80), w=CW, h=Inches(3.80),
              show_data_labels=True, font_name=FONT)

set_footer(slide3)

# ===== SLIDE 4: DARK ROAST (スコアボード) =====
slide4 = prs.slides.add_slide(prs.slide_layouts[2])
clear_placeholders(slide4)

add_breadcrumb(slide4, "Coffee Beans > Dark Roast")
add_title(slide4, "Dark Roast（深煎り）")
add_message_line(slide4, "深い苦味とコク、スモーキーな香り。リッチで複雑な味わい")

# カード型スコアボード（2行×2.5列）
scores_dark = [
    ("酸味", 2),
    ("コク", 5),
    ("甘さ", 2),
    ("苦味", 5),
    ("香り", 5)
]

card_w = Inches(2.0)
card_h = Inches(1.20)
start_x = Inches(0.80)
start_y = Inches(2.00)
gap = Inches(0.30)

for i, (name, score) in enumerate(scores_dark):
    col = i % 5
    row = i // 5

    x = start_x + col * (card_w + gap)
    y = start_y + row * (card_h + gap)

    # カード背景
    shape = slide4.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                     x, y, int(card_w), int(card_h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 245)  # ライトグレー
    shape.line.color.rgb = RGBColor(88, 24, 112)
    shape.line.width = Pt(2)

    # ラベル
    tb = slide4.shapes.add_textbox(x + Inches(0.10), y + Inches(0.10),
                                    int(card_w) - Inches(0.20), Inches(0.40))
    tf = tb.text_frame
    tf.text = name
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(88, 24, 112)
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    # スコア（大きく表示）
    tb2 = slide4.shapes.add_textbox(x + Inches(0.10), y + Inches(0.50),
                                     int(card_w) - Inches(0.20), Inches(0.60))
    tf2 = tb2.text_frame
    tf2.text = f"{score} / 5"
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(88, 24, 112)
    p2.font.name = "Arial"
    p2.alignment = PP_ALIGN.CENTER

set_footer(slide4)

# ===== SAVE =====
output_path = os.path.join(os.getcwd(), "coffee_beans_tasting_card.pptx")
prs.save(output_path)
print(f"✅ Created: {output_path}")
