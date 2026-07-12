import sys, os
_SKILL = os.path.join(os.path.expanduser("~"), ".claude", "skills", "acnpptx", "scripts")
sys.path.insert(0, _SKILL)

import helpers as _h
_h.set_lang("ja")
_h.load_theme("accenture")
from helpers import *

import helpers as _helpers_mod
from pptx.dml.color import RGBColor
_helpers_mod.DARK_PURPLE = RGBColor(0x75, 0x00, 0xC0)
DARK_PURPLE = RGBColor(0x75, 0x00, 0xC0)

from native_shapes import *
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn as _qn

LAYOUT_COVER   = 0
LAYOUT_CONTENT = 2
LAYOUT_SECTION = 6

# Brand-safe purple variants (allowed by brand_check)
P1 = RGBColor(0x46, 0x00, 0x73)  # DARKEST_PURPLE
P2 = RGBColor(0x75, 0x00, 0xC0)  # DARK_PURPLE
P3 = RGBColor(0x75, 0x00, 0xC0)  # same (no mid permitted)
P4 = RGBColor(0x46, 0x00, 0x73)  # DARKEST_PURPLE

prs = Presentation(TEMPLATE_PATH)
while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    prs.part.drop_rel(sldId.get(_qn("r:id")))
    del prs.slides._sldIdLst[0]

def add_slide(layout_idx=LAYOUT_CONTENT):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    clear_placeholders(slide)
    return slide

def _rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def _tb(slide, text, x, y, w, h, size, bold, color, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = FONT; p.alignment = align
    return tb

def _tb_multi(slide, lines, x, y, w, h, size, bold, color, space_after=6, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = color; p.font.name = FONT
        p.space_after = Pt(space_after)
    return tb

# ── SLIDE 1: COVER ────────────────────────────────────────────────
def make_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_COVER])
    # Set slide background explicitly so brand_check reads dark background
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = P1
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        tf = ph.text_frame
        if idx == 0:  # title
            tf.clear()
            p1 = tf.paragraphs[0]
            p1.text = "N1分析レポート"
            p1.font.size = Pt(40); p1.font.bold = True
            p1.font.color.rgb = WHITE; p1.font.name = FONT
            p2 = tf.add_paragraph()
            p2.text = "被相続人 #5 ミラクルピカ（70歳・女性・静岡）"
            p2.font.size = Pt(28); p2.font.bold = True
            p2.font.color.rgb = WHITE; p2.font.name = FONT
        elif idx == 1:  # subtitle
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "インタビュー実施日：2026年7月1日"
            p.font.size = Pt(16); p.font.bold = False
            p.font.color.rgb = WHITE; p.font.name = FONT
        else:
            tf.clear()
            p = tf.paragraphs[0]; p.text = " "; p.font.name = FONT

make_cover()

# ── SLIDE 2: AFFINITY MAP ─────────────────────────────────────────
def make_affinity():
    slide = add_slide()
    add_breadcrumb(slide, "N1分析 > アフィニティマップ")
    add_title(slide, "4つのクラスターから見えるミラクルピカの構造")
    add_message_line(slide, "Factを束ねると「揉め事・任せる・夫婦ズレ・今が窓」の4軸が浮かぶ")
    set_footer(slide)

    # 2×2 grid of colored panels
    clusters = [
        {
            "title": "A｜家族の複雑な力学",
            "color": P1,
            "facts": [
                "揉め事の本質は「長男嫁 vs 次男嫁」",
                "ミラクルピカは「どちらも可愛い」と言いつつ次男を応援",
                "長男と4カ月連絡なし・配分の話が凍結",
                "夫は「全部次男に」の遺言書を書くかもしれない",
            ]
        },
        {
            "title": "B｜多角経営者の行動原理",
            "color": P2,
            "facts": [
                "「面白そう」で動き失敗も笑い話にできる",
                "採算の取れない物件は今年2月に整理済み",
                "毎月複数の収入源を分散構築（プロの発想）",
                "「書かないとダメ」とわかっているが手が動かない",
            ]
        },
        {
            "title": "C｜相談先への徹底的な受動性",
            "color": P2,
            "facts": [
                "レオパレス・証券会社・朝日生命：すべて「来てくれた」",
                "かんぽ保険金、面倒で請求に行かない",
                "無料相談は「無料じゃない入り口が気に食わない」",
                "「向こうから来てくれれば話す」が一貫した姿勢",
            ]
        },
        {
            "title": "D｜先送りと発火のタイムライン",
            "color": P1,
            "facts": [
                "「まだ死なないと思っているから先送り」",
                "半年前からやっと自由な時間ができた",
                "「ほんの何日か前」に朝日生命の友人と保険の話",
                "終活スイッチは今まさに発火しつつある状態",
            ]
        },
    ]

    col_w = CW / 2 - 0.08
    row_h = (BY - CY) / 2 - 0.08
    positions = [
        (ML, CY),
        (ML + col_w + 0.16, CY),
        (ML, CY + row_h + 0.16),
        (ML + col_w + 0.16, CY + row_h + 0.16),
    ]

    for i, (c, (px, py)) in enumerate(zip(clusters, positions)):
        _rect(slide, px, py, col_w, row_h, c["color"])
        # title
        _tb(slide, c["title"], px + 0.12, py + 0.10, col_w - 0.24, 0.32,
            12, True, WHITE, PP_ALIGN.LEFT)
        # facts
        bullet_lines = ["• " + f for f in c["facts"]]
        _tb_multi(slide, bullet_lines, px + 0.12, py + 0.46,
                  col_w - 0.24, row_h - 0.56, 10, False, WHITE, space_after=4)

make_affinity()

# ── SLIDE 3: MINDSET ──────────────────────────────────────────────
def make_mindset():
    slide = add_slide()
    add_breadcrumb(slide, "N1分析 > マインドセット")
    add_title(slide, "ミラクルピカの世界観：「面白そう」で動き「全部任せる」で完結する")
    add_message_line(slide, "相続も「誰かがまるっとやってくれる」前提がそろった時だけ動き出す")
    set_footer(slide)

    # Chevron flow: 3 phases of mindset
    phases = [
        ("行動トリガー", "「面白そう」「任せてくれる人がいる」\nの二条件が揃った時だけ動く"),
        ("相続への姿勢", "揉め事さえなければ幸せ。\n自分の楽しみ優先。でも子供が揉めるくらいなら残さない"),
        ("信頼の構造", "来てくれる・全部やってくれる・\n雑談から始まる——この3条件が信頼の発火条件"),
    ]

    chev_w = (CW - 0.30) / 3
    chev_h = 0.70
    gap = 0.15
    chev_y = CY + 0.10

    for i, (label, desc) in enumerate(phases):
        cx = ML + i * (chev_w + gap)
        # chevron shape
        chev = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON,
            Inches(cx), Inches(chev_y), Inches(chev_w), Inches(chev_h))
        chev.fill.solid(); chev.fill.fore_color.rgb = DARKEST_PURPLE
        chev.line.fill.background()
        # label
        tf = chev.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = label
        p.font.size = Pt(13); p.font.bold = True
        p.font.color.rgb = WHITE; p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER

        # description below
        desc_y = chev_y + chev_h + 0.18
        _tb_multi(slide, desc.split("\n"),
                  cx, desc_y, chev_w, BY - desc_y - 0.05,
                  12, False, TEXT_BODY, space_after=6)

    # Key quote box
    quote_y = chev_y + chev_h + 1.40
    _rect(slide, ML, quote_y, CW, BY - quote_y - 0.05, OFF_WHITE)
    _tb(slide, "「とにかく子供たちが揉めないことが一番の目標。揉めるくらいなら全部使ってしまおう、っていうのがみんなの合言葉になっている」",
        ML + 0.15, quote_y + 0.12, CW - 0.30, BY - quote_y - 0.25,
        12, False, DARKEST_PURPLE, PP_ALIGN.LEFT)

make_mindset()

# ── INSIGHT CARD HELPER ───────────────────────────────────────────
def make_insight(num, title, insight_text, honne, quotes, segment, hanshо, client_meaning, hmw, wow=True):
    slide = add_slide()
    add_breadcrumb(slide, f"N1分析 > Insight #{num}")
    add_title(slide, f"Insight #{num}：{title}")
    msg = "★WOW" if wow else ""
    add_message_line(slide, insight_text[:60] + ("…" if len(insight_text) > 60 else ""))
    set_footer(slide)

    # Left column: insight + honне + quotes
    L_W = CW * 0.55
    R_W = CW - L_W - 0.20
    R_X = ML + L_W + 0.20

    # Insight body
    _rect(slide, ML, CY, L_W, 1.20, DARKEST_PURPLE)
    _tb(slide, insight_text, ML + 0.12, CY + 0.08, L_W - 0.24, 1.04,
        12, False, WHITE, PP_ALIGN.LEFT)

    # Honne
    honne_y = CY + 1.28
    _tb(slide, "背景にある本音", ML, honne_y, L_W, 0.28, 11, True, DARK_PURPLE)
    _rect(slide, ML, honne_y + 0.28, L_W, 0.95, OFF_WHITE)
    _tb(slide, honne, ML + 0.10, honne_y + 0.30, L_W - 0.20, 0.90,
        11, False, TEXT_BODY, PP_ALIGN.LEFT)

    # Quotes
    q_y = honne_y + 1.31
    _tb(slide, "根拠引用（逐語より）", ML, q_y, L_W, 0.28, 11, True, DARK_PURPLE)
    for j, q in enumerate(quotes[:3]):
        qy2 = q_y + 0.28 + j * 0.50
        _rect(slide, ML, qy2, 0.04, 0.38, DARK_PURPLE)
        _tb(slide, f"「{q}」", ML + 0.12, qy2, L_W - 0.12, 0.44,
            10, False, TEXT_BODY, PP_ALIGN.LEFT)

    # Right column: 4 info boxes
    info_items = [
        ("対象セグメント", segment),
        ("反証", hanshо),
        ("クライアントへの示唆", client_meaning),
        ("HMW（問い）", hmw),
    ]
    box_h = (BY - CY - 0.05) / 4 - 0.08
    colors = [P1, P2, P1, P2]
    for j, (label, content) in enumerate(info_items):
        by2 = CY + j * (box_h + 0.08)
        _rect(slide, R_X, by2, R_W, box_h, colors[j])
        _tb(slide, label, R_X + 0.10, by2 + 0.06, R_W - 0.20, 0.24,
            10, True, WHITE, PP_ALIGN.LEFT)
        _tb(slide, content, R_X + 0.10, by2 + 0.28, R_W - 0.20, box_h - 0.34,
            10, False, WHITE, PP_ALIGN.LEFT)

# ── SLIDE 4: INSIGHT #1 ───────────────────────────────────────────
make_insight(
    num=1,
    title="揉め事が相続準備を封印する",
    insight_text="知識も意欲も資産もある。しかし「誰に何を渡すか」が家族間の感情的対立によって決められない状態では、どんな相続準備ツールも機能しない。相続準備の最大の阻害要因は無知でも面倒くさがりでもなく、未解決の家族感情である。",
    honne="「揉めないことが一番」と言いながら、実際は揉め事の渦中にいる。自分では制御できない状況に直面し、「書けない・決められない」という無力感が相続準備を止めている。",
    quotes=[
        "「今なんか自分がちょっとわからないような状態で。どうしたらいいのか」[1:08:46]",
        "「こっちはあなた達が取ればいいっていう話は割にしてたんですけどね。今まだそれがはっきりしていない」[1:08:26]",
        "「書いておかないとダメって思ってるんですけど、まだ」[1:10:14]",
    ],
    segment="事業承継後・子供間に感情的対立がある層（被相続人・資産規模問わず）",
    hanshо="揉め事がなければ相続準備はスムーズに進む可能性がある（こぺぽとの対比）。家族関係が良好な被相続人には機能しない問いかもしれない。",
    client_meaning="相続窓口への誘導より先に「家族関係の整理をサポートする」接点設計が必要。法律や税金の話の前に「感情の交通整理」ができる専門家（ファミリーコーチ・終活カウンセラー）との連携が有効かもしれない。",
    hmw="HMW：「誰に渡すか」が決まっていなくても着手できる相続準備の入り口をどう作れるか？",
    wow=True,
)

# ── SLIDE 5: INSIGHT #2 ───────────────────────────────────────────
make_insight(
    num=2,
    title="信頼は「先に来てくれる人」にしか生まれない",
    insight_text="ミラクルピカの全ての相談・信頼関係は「向こうから来てくれた」という共通構造を持つ。レオパレス・証券会社・朝日生命の友人——全員がアウトリーチ型。逆に、こちらから出向く形の相談（無料窓口・銀行来店）には一度も動いていない。",
    honne="「面倒くさがり屋」というのは性格ではなく、信頼の閾値の問題。「あなたのために来てくれた」という体験が先にあって初めて、相談の扉が開く。まるっと任せられる安心感がないと動かない。",
    quotes=[
        "「積極的に来られると、そういうシステムがあればちょっと相談にっていうことになるかもしれない」[1:48:35]",
        "「向こうから言ってくれればご自宅に伺えますみたいに言っていいですとかって言って」[1:43:35]",
        "「レオパレスが全部やってくれたんですよ…私たちは仕事がすごい忙しかったので」[45:04]",
    ],
    segment="多角経営経験者・高資産層・アナログ重視の70代女性",
    hanshо="証券会社担当との関係は「最初から来てくれた」わけではなく、投資意欲のある人から始まった関係。完全な受動型ではなく、自分の興味と一致した時に自ら関係を深める側面もある。",
    client_meaning="郵便局・ゆうちょの「窓口に来てください」モデルはこの層には機能しない。担当者が自宅・施設へ出向くアウトリーチ型の接点設計が必須。「来てもいいですか？」の一声が最初のアクション。",
    hmw="HMW：ゆうちょの担当者が「先に来る人」として記憶されるにはどんな接点設計が必要か？",
    wow=True,
)

# ── SLIDE 6: INSIGHT #3 ───────────────────────────────────────────
make_insight(
    num=3,
    title="相続の最大リスクは夫婦の中にある",
    insight_text="ミラクルピカは「子供が揉めないように」と心配しているが、より深いリスクは夫婦間の方針のズレにある。「平等に渡したい」ミラクルピカと「全部次男に」という夫。夫が先に遺言書を書けば、その後ミラクルピカが亡くなった時に深刻な問題になりうる。",
    honne="夫への感謝と信頼は深い（「失敗しても文句を言わず後始末をしてくれた」）。だからこそ、夫の方針への異論を直接言いにくく、夫婦間の合意形成が先送りされている。",
    quotes=[
        "「主人はもう長男を見限っていて、全部次男にやるつもりでいます」[1:36:29]",
        "「だから遺言書を書くかもしれない。ちょっとわからないですけどね」[1:39:23]",
        "「主人には怒られるけど、やっぱり長男にもあげたい気持ちもあるし」[33:15前後]",
    ],
    segment="夫婦で資産管理・意思決定を分担してきたシニア層（特に被相続人が複数いるケース）",
    hanshо="夫婦間の方針は話し合えば合意できる可能性もある。「怒られる」という表現が実際の対立の深さを示すかは不明（根拠薄）。",
    client_meaning="相続相談は「個人」ではなく「夫婦単位」で設計する必要がある。夫婦が同席する場（自宅訪問・お茶を飲みながらの雑談形式）でないと片方の本音しか拾えない。夫婦間のズレを可視化するシンプルなツールが有効かもしれない。",
    hmw="HMW：夫婦が「相続の方針」について自然に会話できる場をどう作れるか？",
    wow=False,
)

# ── SLIDE 7: INSIGHT #4 ───────────────────────────────────────────
make_insight(
    num=4,
    title="セカンドライフ直後が唯一の接触窓",
    insight_text="50〜60代は事業・音楽・執筆で時間分刻みの多忙。「来てもいいよ」とも言えなかった。半年前に事業を手放してやっと余白ができた。この「引退直後の時間的余白」と、「ほんの何日か前」に発火した終活スイッチが重なった今が、接触できる最初の窓。窓はごく短い。",
    honne="「まだ死なないから先送り」は言い訳ではなく、本当に今まで自分のことを考える暇がなかった。余白ができた今、やっと「そろそろ考えよう」という気持ちになっている。その気持ちが行動になる前に接触できるかが勝負。",
    quotes=[
        "「やっとここで仕事が少し減ったので楽しみたい、自分達が楽しむことを優先している」[1:23:45]",
        "「本当にここ何日か前、もう本当にすごくゆってもらって。やっぱり保険が一番いいのかなって」[1:24:15]",
        "「銀行に行くっていっても、ほんとにはいはいっていう感じで、余裕が全くなかった」[1:45:07]",
    ],
    segment="事業引退直後・セカンドライフ開始期の60〜70代（特に自営業経験者）",
    hanshо="窓が開いているとはいえ、旅行・趣味に9割のエネルギーが向いている。相続の話が割り込む余地は限られており、「旅行の話からじゃないと入れない」可能性がある。",
    client_meaning="退職・事業移譲タイミングを郵便局・銀行が把握するデータがあれば先手を打てる。「セカンドライフ応援」という文脈（旅行保険・クルーズ特典）での接点が相続相談の入り口になりうる。",
    hmw="HMW：事業引退・退職というタイミングをゆうちょが把握し、自然な形でコンタクトを始めるにはどんな仕組みが必要か？",
    wow=True,
)

# ── SLIDE 8: HMW SUMMARY ─────────────────────────────────────────
def make_hmw():
    slide = add_slide()
    add_breadcrumb(slide, "N1分析 > HMWまとめ")
    add_title(slide, "4つの「どうすれば（HMW）」")
    add_message_line(slide, "ミラクルピカから導かれた設計の問い——いずれも郵政の既存モデルへの挑戦状")
    set_footer(slide)

    hmws = [
        {
            "num": "1",
            "insight": "揉め事が封印",
            "q": "「誰に渡すか」が決まっていなくても着手できる相続準備の入り口をどう作れるか？",
            "color": P1,
        },
        {
            "num": "2",
            "insight": "先に来てくれる人",
            "q": "ゆうちょの担当者が「先に来る人」として記憶されるにはどんな接点設計が必要か？",
            "color": P2,
        },
        {
            "num": "3",
            "insight": "夫婦間のズレ",
            "q": "夫婦が「相続の方針」について自然に会話できる場をどう作れるか？",
            "color": P1,
        },
        {
            "num": "4",
            "insight": "セカンドライフ直後の窓",
            "q": "事業引退・退職タイミングをゆうちょが把握し、自然な形でコンタクトを始めるにはどんな仕組みが必要か？",
            "color": P2,
        },
    ]

    card_w = (CW - 0.30) / 4
    card_h = BY - CY - 0.05

    for i, item in enumerate(hmws):
        cx = ML + i * (card_w + 0.10)
        _rect(slide, cx, CY, card_w, card_h, item["color"])
        # number badge: WHITE text on card background
        _tb(slide, item["num"], cx + 0.10, CY + 0.08, 0.32, 0.36,
            20, True, WHITE, PP_ALIGN.LEFT, wrap=False)
        # insight label
        _tb(slide, f"Insight #{item['num']}\n{item['insight']}",
            cx + 0.10, CY + 0.50, card_w - 0.20, 0.60,
            10, True, WHITE, PP_ALIGN.LEFT)
        # HMW label: white-bordered box approach
        _tb(slide, "▶ HMW", cx + 0.10, CY + 1.18, card_w - 0.20, 0.28,
            10, True, WHITE, PP_ALIGN.CENTER, wrap=False)
        # Question text
        _tb(slide, item["q"], cx + 0.10, CY + 1.52,
            card_w - 0.20, card_h - 1.60,
            11, False, WHITE, PP_ALIGN.LEFT)

# ── CLOSING ───────────────────────────────────────────────────────
make_closing_slide(prs, text_color=BLACK)
strip_sections(prs)

out = os.path.join(os.path.expanduser("~"), "Documents", "Obsidian Vault", "ミラクルピカ_N1分析.pptx")
prs.save(out)
print(f"Saved: {out}")
