import sys, os
_SKILL = os.path.join(os.path.expanduser("~"), ".claude", "skills", "acnpptx", "scripts")
sys.path.insert(0, _SKILL)

import helpers as _h
_h.set_lang("ja")
_h.load_theme("acn-slide-master")
from helpers import *

from native_shapes import *
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn as _qn

prs = Presentation(TEMPLATE_PATH)
while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    prs.part.drop_rel(sldId.get(_qn("r:id")))
    del prs.slides._sldIdLst[0]

LAYOUT_COVER = 0
LAYOUT_CONTENT = 57


def add_slide(layout_idx=LAYOUT_CONTENT):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    clear_placeholders(slide)
    return slide


def add_cover_slide():
    return prs.slides.add_slide(prs.slide_layouts[LAYOUT_COVER])


# ─── Slide 1: Cover ─────────────────────────────────────────────────────────
slide = add_cover_slide()
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    ph.text_frame.clear()
    p = ph.text_frame.paragraphs[0]
    if idx == 0:
        p.text = "終活サービス"
        p.font.size = Pt(40); p.font.bold = True
        p.font.color.rgb = WHITE; p.font.name = FONT
    elif idx == 1:
        p.text = "人生の最終章を、安心・丁寧にサポートする"
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE; p.font.name = FONT
    elif idx == 12:
        p.text = "Accenture"
        p.font.size = Pt(16); p.font.color.rgb = LIGHT_PURPLE; p.font.name = FONT
    elif idx == 15:
        p.text = "2026年6月"
        p.font.size = Pt(8); p.font.color.rgb = LIGHT_PURPLE; p.font.name = FONT
    else:
        p.text = " "
        p.font.size = Pt(8); p.font.name = FONT


# ─── Slide 2: Agenda (Pattern I) ────────────────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "アジェンダ")
add_title(slide, "アジェンダ")

items = [
    "終活とは？",
    "市場概況",
    "主要サービスカテゴリ",
    "ターゲット顧客",
    "サービス提供モデル",
    "競合・差別化",
    "今後の展望",
]
ITEM_Y = CY + 0.10
gap = min(0.78, (BY - ITEM_Y - 0.65) / max(len(items) - 1, 1))

for i, item in enumerate(items):
    iy = ITEM_Y + i * gap
    num_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(ML), Inches(iy), Inches(0.55), Inches(0.62))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = DARKEST_PURPLE
    num_box.line.fill.background()
    tf = num_box.text_frame
    tf.word_wrap = False
    _bp = tf._txBody.find(_qn("a:bodyPr"))
    for _attr in ("lIns", "rIns", "tIns", "bIns"):
        _bp.set(_attr, "0")
    _bp.set("anchor", "ctr")
    tf.paragraphs[0].text = f"{i+1:02d}"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = FONT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(
        Inches(ML + 0.70), Inches(iy + 0.08), Inches(11.30), Inches(0.50))
    p = tb.text_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_BODY
    p.font.name = FONT

set_footer(slide)


# ─── Slide 3: 終活とは何か (Pattern P) ──────────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "終活の概要")
add_title(slide, "終活とは何か")
add_message_line(slide, "終活は死を準備するのではなく、今をより豊かに生きるための活動である")

steps = ["定義と歴史", "社会的背景", "意識の変化", "今後の方向性"]
details = [
    "• 2009年「週刊朝日」が命名\n• 人生の終末に向けた前向きな準備活動\n• 遺言・相続・葬儀の事前整理",
    "• 高齢化率29%（2024年、世界最高水準）\n• 単身高齢者世帯：約700万世帯\n• 看取り難民・孤独死の社会問題化",
    "• 内閣府調査：60代の68%が終活を「必要」と回答\n• 40〜50代の若年化：早期着手が主流に\n• タブー感の薄れ、前向きな活動へ転換",
    "• デジタル終活の急拡大（SNS・暗号資産）\n• AIによるパーソナル終活プランの自動生成\n• 地方自治体との連携サービスが増加",
]

FLOW_H = 0.80
GAP = 0.25
DETAIL_H = 2.80
FLOW_Y = CY + (AH - FLOW_H - GAP - DETAIL_H) / 2
DETAIL_Y = FLOW_Y + FLOW_H + GAP

add_chevron_flow(
    slide, steps,
    x=ML, y=FLOW_Y, total_w=CW, h=FLOW_H,
    fill_color=DARKEST_PURPLE,
    text_color=WHITE,
    shape_style="chevron",
    use_pentagon_first=True,
    font_size_pt=14,
    font_name=FONT,
)

n = len(steps)
col_w = CW / n
for i, detail in enumerate(details):
    dx = ML + i * col_w
    tb = slide.shapes.add_textbox(Inches(dx + 0.10), Inches(DETAIL_Y),
        Inches(col_w - 0.20), Inches(DETAIL_H))
    tf = tb.text_frame; tf.word_wrap = True
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "t")
    for j, line in enumerate(detail.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(14)
        p.font.color.rgb = TEXT_BODY; p.font.name = FONT
        p.space_after = Pt(8)

set_footer(slide)


# ─── Slide 4: 市場概況 (Pattern W) ──────────────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "市場概況")
add_title(slide, "高齢化社会と市場規模")
add_message_line(slide, "2030年に向け終活関連市場は1.5兆円超に拡大、参入機会は今が最大")

stats = [
    {"value": "29%", "label": "高齢化率（2024年）",
     "detail": "世界最高水準。2040年には35%超へ。終活需要の根本的なドライバー",
     "source": "出典: 総務省統計局 人口推計（2024年）"},
    {"value": "1.5兆円", "label": "市場規模（2030年予測）",
     "detail": "終活関連市場全体。2024年比約2倍。葬儀・相続・介護連携が牽引",
     "source": "出典: 矢野経済研究所「シニア市場白書2024」"},
    {"value": "8%", "label": "年平均成長率（CAGR）",
     "detail": "デジタル終活・事前予約型葬儀など新サービスが市場拡大を加速",
     "source": "出典: 富士経済「終活・シニアサービス市場分析2024」"},
    {"value": "2,500万人", "label": "潜在顧客数（65歳以上）",
     "detail": "60代以上のうち終活未着手者は約8割。潜在的な獲得機会は巨大",
     "source": "出典: 内閣府「高齢社会白書2024」"},
]

n = len(stats)
col_w = CW / n

for i in range(n - 1):
    div_x = ML + (i + 1) * col_w
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(div_x - 0.01), Inches(CY), Inches(0.02), Inches(BY - CY - 0.20))
    line.fill.solid(); line.fill.fore_color.rgb = LIGHT_GRAY
    line.line.fill.background()

for i, stat in enumerate(stats):
    cx = ML + i * col_w
    cw = col_w - 0.10

    tb_val = slide.shapes.add_textbox(Inches(cx + 0.05), Inches(CY + 0.20), Inches(cw), Inches(1.20))
    tb_val.text_frame.word_wrap = True
    p = tb_val.text_frame.paragraphs[0]
    p.text = stat["value"]
    p.font.size = Pt(44); p.font.bold = True
    p.font.color.rgb = DARKEST_PURPLE; p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    tb_lbl = slide.shapes.add_textbox(Inches(cx + 0.05), Inches(CY + 1.45), Inches(cw), Inches(0.55))
    p = tb_lbl.text_frame.paragraphs[0]
    p.text = stat["label"]
    p.font.size = Pt(14); p.font.bold = True
    p.font.color.rgb = TEXT_BODY; p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    tb_det = slide.shapes.add_textbox(Inches(cx + 0.05), Inches(CY + 2.10), Inches(cw), Inches(1.80))
    tf = tb_det.text_frame; tf.word_wrap = True
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "t")
    p = tf.paragraphs[0]
    p.text = stat["detail"]
    p.font.size = Pt(14); p.font.color.rgb = TEXT_BODY; p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(6)

    tb_src = slide.shapes.add_textbox(Inches(cx + 0.05), Inches(BY - 0.60), Inches(cw), Inches(0.55))
    tb_src.text_frame.word_wrap = True
    p = tb_src.text_frame.paragraphs[0]
    p.text = stat["source"]
    p.font.size = Pt(12); p.font.color.rgb = MID_GRAY; p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

set_footer(slide)


# ─── Slide 5: 主要サービスカテゴリ (Pattern B) ─────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "サービス詳細")
add_title(slide, "主要サービスカテゴリ")
add_message_line(slide, "4つの主要カテゴリが終活の全プロセスをカバーし、相互補完的に機能する")

panels = [
    {
        "title": "法的・財務サポート",
        "items": [
            "• 遺言書作成・公正証書化支援（弁護士・司法書士連携）",
            "• 相続手続き代行・相続税対策コンサルティング",
            "• 財産目録作成・資産整理（不動産・金融資産・デジタル資産）",
            "• 成年後見制度の活用支援・任意後見契約の締結サポート",
        ]
    },
    {
        "title": "ライフエンド・葬儀サポート",
        "items": [
            "• 葬儀・告別式の事前予約（全国提携500社以上）",
            "• 墓地・納骨堂の選定・契約代行（樹木葬・散骨も対応）",
            "• エンディングノートの作成支援とデジタル保管",
            "• 遺品整理・形見分け・不要品処分の一括請負",
        ]
    },
    {
        "title": "デジタル終活",
        "items": [
            "• SNS・メールアカウントの死後処理代行（Facebook追悼設定等）",
            "• 暗号資産・デジタル金融資産の引き継ぎ設計",
            "• パスワード・ID情報の安全な信託保管と相続",
            "• デジタル遺書・動画メッセージの作成・保管サービス",
        ]
    },
    {
        "title": "医療・介護連携",
        "items": [
            "• ACP（人生会議）支援・尊厳死の意思表示書類作成",
            "• かかりつけ医・介護事業者との多職種連携プラットフォーム",
            "• 介護保険の申請代行・認定調査への同席支援",
            "• ホスピス・在宅緩和ケアの選定・コーディネーション",
        ]
    },
]

# 2×2グリッド配置
CARD_W = (CW - 0.20) / 2
CARD_H = (BY - CY - 0.50 - 0.20) / 2
xs = [ML, ML + CARD_W + 0.20]
ys_pos = [CY + 0.30, CY + 0.30 + CARD_H + 0.20]

for idx_p, panel in enumerate(panels):
    px = xs[idx_p % 2]
    py = ys_pos[idx_p // 2]

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(px), Inches(py), Inches(CARD_W), Inches(CARD_H))
    bg.fill.solid(); bg.fill.fore_color.rgb = OFF_WHITE; bg.line.fill.background()

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(px), Inches(py), Inches(CARD_W), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = CORE_PURPLE; bar.line.fill.background()

    tb_h = slide.shapes.add_textbox(Inches(px + 0.15), Inches(py + 0.14),
        Inches(CARD_W - 0.30), Inches(0.40))
    p = tb_h.text_frame.paragraphs[0]
    p.text = panel["title"]
    p.font.size = Pt(14); p.font.bold = True
    p.font.color.rgb = BLACK; p.font.name = FONT

    body_y = py + 0.60
    body_h = CARD_H - 0.70
    tb_b = slide.shapes.add_textbox(Inches(px + 0.15), Inches(body_y),
        Inches(CARD_W - 0.30), Inches(body_h))
    tf = tb_b.text_frame; tf.word_wrap = True
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "t")
    for j, line in enumerate(panel["items"]):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(12)
        p.font.color.rgb = TEXT_BODY; p.font.name = FONT
        p.space_after = Pt(6)

set_footer(slide)


# ─── Slide 6: ターゲット顧客 (Pattern K) ────────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "ターゲット顧客")
add_title(slide, "ターゲット顧客セグメント")
add_message_line(slide, "60代以上のシニア層を主軸に、家族・法人の3層構造でアプローチする")

N_COLS = 3
PANEL_W = (CW - 0.20 * (N_COLS - 1)) / N_COLS
PANEL_H = BY - CY - 0.60
BADGE_D = 0.48

col_data = [
    {
        "num": 1, "title": "シニア層（60〜80代）",
        "body": "【規模】約2,500万人（65歳以上）\n\n• 資産保有額1,000万円以上の富裕層シニア\n• 子どもと離れて暮らす単身・夫婦世帯\n• 持病・介護不安を抱える70代以上\n• 平均LTV：月額5,800円×36ヶ月=20.9万円"
    },
    {
        "num": 2, "title": "家族・介護者（40〜60代）",
        "body": "【規模】約1,200万人（要介護者の主介護者）\n\n• 親の終活を遠方から支援したい子世代\n• ダブルケア（育児×介護）世帯の時間節約ニーズ\n• 突然の相続・葬儀で困った経験のある層\n• 平均LTV：月額2,980円×24ヶ月=7.2万円"
    },
    {
        "num": 3, "title": "企業・法人",
        "body": "【規模】約85万社（従業員1,000人以上）\n\n• 福利厚生として終活サービスを導入する大企業\n• 生命保険・損害保険会社との白ラベル提携\n• 自治体・病院・介護施設との公共連携契約\n• 平均契約額：年間300万円×3年=900万円"
    },
]

for i, col in enumerate(col_data):
    px = ML + i * (PANEL_W + 0.20)
    py = CY + 0.55

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(px), Inches(py), Inches(PANEL_W), Inches(PANEL_H))
    bg.fill.solid(); bg.fill.fore_color.rgb = OFF_WHITE; bg.line.fill.background()

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(px), Inches(py), Inches(PANEL_W), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = CORE_PURPLE; bar.line.fill.background()

    bx = px + (PANEL_W - BADGE_D) / 2
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(bx), Inches(py - BADGE_D / 2), Inches(BADGE_D), Inches(BADGE_D))
    badge.fill.solid(); badge.fill.fore_color.rgb = DARKEST_PURPLE; badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = str(col["num"])
    p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "ctr")

    tb_t = slide.shapes.add_textbox(Inches(px + 0.12), Inches(py + 0.18),
        Inches(PANEL_W - 0.24), Inches(0.50))
    p = tb_t.text_frame.paragraphs[0]
    p.text = col["title"]
    p.font.size = Pt(14); p.font.bold = True
    p.font.color.rgb = BLACK; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER

    tb_b = slide.shapes.add_textbox(Inches(px + 0.12), Inches(py + 0.75),
        Inches(PANEL_W - 0.24), Inches(PANEL_H - 0.90))
    tf = tb_b.text_frame; tf.word_wrap = True
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "t")
    for j, line in enumerate(col["body"].split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(12)
        p.font.color.rgb = TEXT_BODY; p.font.name = FONT
        p.space_after = Pt(5)

set_footer(slide)


# ─── Slide 7: サービス提供モデル (Pattern T) ────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "サービス提供モデル")
add_title(slide, "サービス提供モデルと収益構造")
add_message_line(slide, "月額サブスクと成果報酬の組み合わせで安定収益と高LTVを両立する")

LABEL_W = 1.80
ARROW_H = 0.45
GAP = 0.15
section_h = (BY - CY - 0.10 - ARROW_H - GAP * 3) / 2

sections_data = [
    {
        "label": "課題\n（現状）",
        "items": [
            "• 複数業者への分散：葬儀・相続・介護をそれぞれ別業者に依頼、調整コスト大",
            "• 高額・不透明なコスト：葬儀費用の全国平均195万円、相続手続き費用50〜100万円",
            "• 情報管理の複雑さ：遺言書・保険証書・デジタルIDが分散し家族が把握困難",
        ]
    },
    {
        "label": "ソリューション\n（当社）",
        "items": [
            "• ワンストップサービス：終活全プロセスを一社で完結。専任コンシェルジュが伴走",
            "• 月額980円〜の明確料金：ベーシック/スタンダード/プレミアムの3プランで透明化",
            "• デジタル一元管理：専用アプリで遺言・保険・デジタル資産を安全に集中管理",
        ]
    },
]

for i, sec in enumerate(sections_data):
    sy = CY + 0.05 + i * (section_h + ARROW_H + GAP)

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(ML), Inches(sy), Inches(CW), Inches(section_h))
    bg.fill.solid(); bg.fill.fore_color.rgb = OFF_WHITE
    bg.line.color.rgb = LIGHT_GRAY; bg.line.width = Pt(0.75)

    lbl = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(ML), Inches(sy), Inches(LABEL_W), Inches(section_h))
    lbl.fill.solid(); lbl.fill.fore_color.rgb = DARKEST_PURPLE; lbl.line.fill.background()
    tf = lbl.text_frame; tf.word_wrap = True
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "ctr")
    p = tf.paragraphs[0]; p.text = sec["label"]
    p.font.size = Pt(14); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER

    content_x = ML + LABEL_W + 0.15
    content_w = CW - LABEL_W - 0.20
    tb = slide.shapes.add_textbox(Inches(content_x), Inches(sy + 0.12),
        Inches(content_w), Inches(section_h - 0.25))
    tf = tb.text_frame; tf.word_wrap = True
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "ctr")
    for j, line in enumerate(sec["items"]):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(14)
        p.font.color.rgb = TEXT_BODY; p.font.name = FONT
        p.space_after = Pt(10)

    if i == 0:
        arrow_y = sy + section_h + GAP
        arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
            Inches(ML + CW / 2 - 0.50), Inches(arrow_y), Inches(1.0), Inches(ARROW_H))
        arr.fill.solid(); arr.fill.fore_color.rgb = DARKEST_PURPLE; arr.line.fill.background()

set_footer(slide)


# ─── Slide 8: 競合分析 (Pattern G) ──────────────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "競合・差別化")
add_title(slide, "競合分析と差別化ポイント")
add_message_line(slide, "専門性×デジタル×ワンストップの三位一体で既存競合に対し構造的優位を持つ")

headers = ["評価軸", "当社", "競合A（葬儀特化）", "競合B（法律特化）", "競合C（介護特化）"]
rows_data = [
    ["専門性・カバレッジ", "◎ 全領域ワンストップ", "△ 葬儀・墓のみ", "△ 相続・遺言のみ", "△ 介護・医療のみ"],
    ["デジタル対応", "◎ 専用アプリ・AI提案", "× デジタルなし", "○ 書類電子化のみ", "△ 介護記録のみ"],
    ["価格透明性", "◎ 月額980円〜固定", "△ 都度見積もり高額", "△ 時間制100,000円/h", "○ 介護保険適用"],
    ["全国対応", "◎ 全国47都道府県", "○ 主要30都市", "○ 主要20都市", "△ 地域限定"],
    ["パートナー連携", "◎ 1,000社以上", "△ 独自ネットワーク", "△ 法律専門家のみ", "△ 介護施設のみ"],
    ["コンシェルジュ対応", "◎ 専任担当制24h", "△ 繁忙期のみ", "○ 予約制対応", "○ ケアマネ連携"],
]

tbl_rows = len(rows_data) + 1
tbl_cols = len(headers)
table = slide.shapes.add_table(
    tbl_rows, tbl_cols, Inches(ML), Inches(2.10), Inches(CW), Inches(BY - 2.20)).table

col_widths = [2.20, 2.35, 2.35, 2.35, 3.25]
for c, w in enumerate(col_widths):
    table.columns[c].width = Inches(w)

for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = DARKEST_PURPLE
    for para in cell.text_frame.paragraphs:
        para.font.bold = True; para.font.size = Pt(12)
        para.font.color.rgb = WHITE; para.font.name = FONT
        para.alignment = PP_ALIGN.CENTER

for r, row in enumerate(rows_data, 1):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(12)
            if c == 1:
                para.font.bold = True
                para.font.color.rgb = CORE_PURPLE
            else:
                para.font.color.rgb = TEXT_BODY
            para.font.name = FONT
            para.alignment = PP_ALIGN.CENTER

set_footer(slide)


# ─── Slide 9: 今後の展望 (Pattern E) ────────────────────────────────────────
slide = add_slide()
add_breadcrumb(slide, "今後の展望")
add_title(slide, "今後の展望とロードマップ")
add_message_line(slide, "2027年までに全国展開・提携1,000社・ARR50億円を達成する")

roadmap_items = [
    (
        "2025年（FY25）: 都市圏パイロット展開",
        "東京・大阪・名古屋の3大都市圏でサービス開始。月間500件の終活相談を達成目標。"
        "提携先150社・ARR3億円・NPS+50以上を確立する。主要チャネルはデジタル広告とシニア向け金融機関連携。"
    ),
    (
        "2026年（FY26）: 主要10都市展開・提携500社",
        "全国主要10都市に拡大。提携先500社・月間2,500件・ARR15億円突破を目指す。"
        "自治体向けBtoBサービスを本格化。生命保険会社3社との白ラベル提携を開始し、流入チャネルを多様化する。"
    ),
    (
        "2027年（FY27）: 全国展開・ARR50億円・上場準備",
        "全国47都道府県にサービス展開。提携1,000社・月間10,000件・ARR50億円を達成。"
        "AIによる自動終活プランニング機能をリリース。東証グロース市場への上場準備を開始する。"
    ),
]

row_h = (BY - CY - 0.40) / len(roadmap_items)
for i, (headline, detail) in enumerate(roadmap_items):
    row_y = CY + 0.20 + i * row_h

    acc = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(ML), Inches(row_y + 0.05), Inches(0.06), Inches(0.28))
    acc.fill.solid(); acc.fill.fore_color.rgb = CORE_PURPLE; acc.line.fill.background()

    tb_h = slide.shapes.add_textbox(Inches(ML + 0.18), Inches(row_y),
        Inches(CW - 0.18), Inches(0.40))
    p = tb_h.text_frame.paragraphs[0]
    p.text = headline; p.font.size = Pt(18)
    p.font.bold = True; p.font.color.rgb = TEXT_BODY; p.font.name = FONT

    tb_d = slide.shapes.add_textbox(Inches(ML + 0.18), Inches(row_y + 0.44),
        Inches(CW - 0.18), Inches(row_h - 0.55))
    tf = tb_d.text_frame; tf.word_wrap = True
    tf._txBody.find(_qn("a:bodyPr")).set("anchor", "t")
    p = tf.paragraphs[0]
    p.text = detail; p.font.size = Pt(14)
    p.font.color.rgb = TEXT_BODY; p.font.name = FONT
    p.space_after = Pt(8)

set_footer(slide)


# ─── Closing ─────────────────────────────────────────────────────────────────
make_closing_slide(prs)
strip_sections(prs)

output_path = os.path.join(os.getcwd(), "終活サービス.pptx")
prs.save(output_path)
print(f"✅ 保存完了: {output_path}")
